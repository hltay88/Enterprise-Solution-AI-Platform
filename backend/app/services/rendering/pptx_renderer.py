"""PPTX rendering for presentation deliverables (ATLAS-044)."""

from __future__ import annotations

from io import BytesIO
from typing import Any


def render_presentation_pptx(
    *,
    title: str,
    status: str,
    slides: list[dict[str, Any]],
) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Widescreen-ish default
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for slide_data in slides:
        slide = prs.slides.add_slide(blank)
        slide_title = str(slide_data.get("title") or "Slide")
        key_message = ""
        body = ""
        notes = ""
        review_required = False
        for item in slide_data.get("content_items") or []:
            structured = item.get("structured_data") or {}
            slide_meta = structured.get("slide") or {}
            if slide_meta.get("key_message"):
                key_message = str(slide_meta.get("key_message") or "")
            if item.get("content_type") == "speaker_notes":
                notes = str(item.get("text") or "")
            elif not body:
                body = str(item.get("text") or "")
            if item.get("review_required"):
                review_required = True
            if slide_meta.get("speaker_notes") and not notes:
                notes = str(slide_meta.get("speaker_notes") or "")

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(32)
        p.font.bold = True

        msg_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(1.2))
        msg_frame = msg_box.text_frame
        msg_frame.word_wrap = True
        mp = msg_frame.paragraphs[0]
        prefix = "[REVIEW REQUIRED] " if review_required else ""
        mp.text = prefix + (key_message or body or "")
        mp.font.size = Pt(22)

        body_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12), Inches(3.5))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        bp = body_frame.paragraphs[0]
        bp.text = body if key_message else ""
        bp.font.size = Pt(16)

        if str(status).lower() != "approved":
            draft_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12), Inches(0.4))
            dp = draft_box.text_frame.paragraphs[0]
            dp.text = "DRAFT — NOT APPROVED FOR CUSTOMER RELEASE"
            dp.font.size = Pt(12)
            dp.font.italic = True

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Ensure at least one slide
    if not slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        box.text_frame.paragraphs[0].text = title or "Presentation"

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
