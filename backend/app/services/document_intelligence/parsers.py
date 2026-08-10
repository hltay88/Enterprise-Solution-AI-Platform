"""Native parsers + OCR-backed extractors for Phase 2.1 file types."""

from __future__ import annotations

import csv
import logging
import re
import shutil
import subprocess
from io import StringIO
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader

from app.core.exceptions import ValidationAppError
from app.services.document_intelligence.normalize import detect_language_hint, normalize_text, word_count
from app.services.document_intelligence.ocr import (
    MANUAL_REVIEW_CONFIDENCE,
    OCR_ENGINE_NAME,
    ocr_image,
    ocr_pdf_pages,
    tesseract_available,
)
from app.services.document_intelligence.types import ExtractedPage, ExtractionResult

logger = logging.getLogger(__name__)

# If native PDF text is thinner than this across the file, try OCR.
_PDF_NATIVE_MIN_CHARS = 40


def extract_document(file_path: Path, file_type: str) -> ExtractionResult:
    file_type = file_type.lower()
    if file_type == "pdf":
        return _extract_pdf(file_path)
    if file_type == "docx":
        return _extract_docx(file_path)
    if file_type == "doc":
        return _extract_doc(file_path)
    if file_type == "xlsx":
        return _extract_xlsx(file_path)
    if file_type == "csv":
        return _extract_csv(file_path)
    if file_type == "txt":
        return _extract_txt(file_path)
    if file_type in {"md", "markdown"}:
        return _extract_markdown(file_path)
    if file_type == "pptx":
        return _extract_pptx(file_path)
    if file_type in {"png", "jpg", "jpeg"}:
        return _extract_image(file_path)
    raise ValidationAppError(f"Unsupported file type: {file_type}")


def _finalize(
    pages: list[ExtractedPage],
    *,
    ocr_used: bool,
    metadata: dict[str, str] | None = None,
    warnings: list[str] | None = None,
) -> ExtractionResult:
    normalized_pages: list[ExtractedPage] = []
    needs_review = False
    for page in pages:
        text = normalize_text(page.text)
        conf = page.confidence
        if conf is not None and conf < MANUAL_REVIEW_CONFIDENCE:
            needs_review = True
        lang = page.language or detect_language_hint(text)
        normalized_pages.append(
            ExtractedPage(
                page_number=page.page_number,
                text=text,
                confidence=conf,
                ocr_engine=page.ocr_engine,
                processing_ms=page.processing_ms,
                language=lang,
            ),
        )

    full_text = "\n\n".join(p.text for p in normalized_pages if p.text).strip()
    language = detect_language_hint(full_text)
    meta = dict(metadata or {})
    meta.setdefault("page_count", str(len(normalized_pages)))
    meta.setdefault("char_count", str(len(full_text)))
    meta.setdefault("word_count", str(word_count(full_text)))
    meta.setdefault("ocr_used", "true" if ocr_used else "false")
    if language:
        meta.setdefault("language", language)

    return ExtractionResult(
        pages=normalized_pages,
        full_text=full_text,
        ocr_used=ocr_used,
        language=language,
        needs_manual_review=needs_review,
        metadata=meta,
        warnings=list(warnings or []),
    )


def _extract_pdf(file_path: Path) -> ExtractionResult:
    reader = PdfReader(str(file_path))
    native_pages: list[ExtractedPage] = []
    native_chars = 0
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        native_chars += len(text.strip())
        native_pages.append(ExtractedPage(page_number=index, text=text))

    meta = {
        "parser": "pypdf",
        "pdf_page_count": str(len(reader.pages)),
    }
    if reader.metadata:
        title = getattr(reader.metadata, "title", None)
        author = getattr(reader.metadata, "author", None)
        if title:
            meta["title"] = str(title)
        if author:
            meta["author"] = str(author)

    warnings: list[str] = []
    if native_chars >= _PDF_NATIVE_MIN_CHARS:
        return _finalize(native_pages, ocr_used=False, metadata=meta, warnings=warnings)

    # Scanned / empty PDF → OCR
    if not tesseract_available():
        warnings.append("PDF has little native text and OCR is unavailable")
        return _finalize(native_pages, ocr_used=False, metadata=meta, warnings=warnings)

    try:
        ocr_results = ocr_pdf_pages(file_path)
    except Exception as exc:
        logger.warning("PDF OCR failed for %s: %s", file_path.name, exc)
        warnings.append(f"PDF OCR failed: {exc}")
        return _finalize(native_pages, ocr_used=False, metadata=meta, warnings=warnings)

    pages = [
        ExtractedPage(
            page_number=page_number,
            text=text,
            confidence=confidence,
            ocr_engine=OCR_ENGINE_NAME,
            processing_ms=processing_ms,
        )
        for page_number, text, confidence, processing_ms in ocr_results
    ]
    meta["parser"] = "pypdf+tesseract"
    return _finalize(pages, ocr_used=True, metadata=meta, warnings=warnings)


def _extract_docx(file_path: Path) -> ExtractionResult:
    document = Document(str(file_path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text and paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    core = document.core_properties
    meta = {"parser": "python-docx"}
    if core.title:
        meta["title"] = str(core.title)
    if core.author:
        meta["author"] = str(core.author)
    pages = [ExtractedPage(page_number=1, text=text)]
    return _finalize(pages, ocr_used=False, metadata=meta)


def _extract_doc(file_path: Path) -> ExtractionResult:
    warnings: list[str] = []
    text = ""
    if shutil.which("antiword"):
        try:
            completed = subprocess.run(
                ["antiword", "-w", "0", str(file_path)],
                check=False,
                capture_output=True,
                text=True,
                timeout=60,
            )
            if completed.returncode == 0:
                text = completed.stdout or ""
            else:
                warnings.append(completed.stderr.strip() or "antiword failed")
        except Exception as exc:
            warnings.append(f"antiword error: {exc}")
    else:
        warnings.append("antiword not installed; using binary text fallback for .doc")

    if not text.strip():
        text = _binary_text_fallback(file_path)

    if not text.strip():
        raise ValidationAppError("Unable to extract text from DOC file")

    pages = [ExtractedPage(page_number=1, text=text)]
    return _finalize(pages, ocr_used=False, metadata={"parser": "antiword-or-fallback"}, warnings=warnings)


def _binary_text_fallback(file_path: Path) -> str:
    raw = file_path.read_bytes()
    # Prefer UTF-16LE runs common in Word .doc streams, then latin-1 printable.
    candidates: list[str] = []
    try:
        utf16 = raw.decode("utf-16le", errors="ignore")
        candidates.append(utf16)
    except Exception:
        pass
    candidates.append(raw.decode("latin-1", errors="ignore"))

    best = ""
    for candidate in candidates:
        # Keep runs of printable characters.
        pieces = re.findall(r"[\t\n\r\x20-\x7e\u00a0-\u024f]{4,}", candidate)
        joined = "\n".join(p.strip() for p in pieces if p.strip())
        if len(joined) > len(best):
            best = joined
    return best


def _extract_xlsx(file_path: Path) -> ExtractionResult:
    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    parts: list[str] = []
    sheet_count = 0
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_count += 1
        parts.append(f"## Sheet: {sheet_name}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if cell is None else str(cell).strip() for cell in row]
            if any(cells):
                parts.append("\t".join(cells))
    workbook.close()
    text = "\n".join(parts)
    pages = [ExtractedPage(page_number=1, text=text)]
    return _finalize(
        pages,
        ocr_used=False,
        metadata={"parser": "openpyxl", "sheet_count": str(sheet_count)},
    )


def _extract_csv(file_path: Path) -> ExtractionResult:
    raw = file_path.read_bytes()
    decoded = None
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            decoded = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if decoded is None:
        raise ValidationAppError("Unable to decode CSV file")

    reader = csv.reader(StringIO(decoded))
    rows = ["\t".join(cell.strip() for cell in row) for row in reader if any(cell.strip() for cell in row)]
    text = "\n".join(rows)
    pages = [ExtractedPage(page_number=1, text=text)]
    return _finalize(pages, ocr_used=False, metadata={"parser": "csv", "row_count": str(len(rows))})


def _extract_txt(file_path: Path) -> ExtractionResult:
    raw = file_path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = raw.decode(encoding)
            pages = [ExtractedPage(page_number=1, text=text)]
            return _finalize(pages, ocr_used=False, metadata={"parser": "txt", "encoding": encoding})
        except UnicodeDecodeError:
            continue
    raise ValidationAppError("Unable to decode TXT file")


def _extract_markdown(file_path: Path) -> ExtractionResult:
    """Markdown treated as UTF text with heading hints for section provenance."""
    result = _extract_txt(file_path)
    result.metadata["parser"] = "markdown"
    headings = [
        line.lstrip("#").strip()
        for line in (result.full_text or "").splitlines()
        if line.startswith("#")
    ]
    if headings:
        result.metadata["heading_count"] = str(len(headings))
        # Stash first few headings for callers that read metadata.
        for index, heading in enumerate(headings[:20], start=1):
            result.metadata[f"heading_{index}"] = heading
    return result


def _extract_pptx(file_path: Path) -> ExtractionResult:
    """Extract slide text via python-pptx (reuse output stack for ingest)."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValidationAppError("PPTX support requires python-pptx") from exc

    presentation = Presentation(str(file_path))
    pages: list[ExtractedPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)
        pages.append(
            ExtractedPage(
                page_number=index,
                text="\n".join(parts),
            ),
        )
    if not pages:
        pages = [ExtractedPage(page_number=1, text="")]
    return _finalize(
        pages,
        ocr_used=False,
        metadata={"parser": "python-pptx", "slide_count": str(len(pages))},
    )


def _extract_image(file_path: Path) -> ExtractionResult:
    if not tesseract_available():
        raise ValidationAppError(
            "Image OCR requires tesseract-ocr. Rebuild the backend image with OCR system packages.",
        )
    try:
        text, confidence, processing_ms = ocr_image(file_path)
    except Exception as exc:
        raise ValidationAppError(f"Image OCR failed: {exc}") from exc

    pages = [
        ExtractedPage(
            page_number=1,
            text=text,
            confidence=confidence,
            ocr_engine=OCR_ENGINE_NAME,
            processing_ms=processing_ms,
        ),
    ]
    return _finalize(
        pages,
        ocr_used=True,
        metadata={"parser": "tesseract"},
    )
