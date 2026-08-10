"""Sprint 5.1 — knowledge file ingest parsers (PPTX, Markdown)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.constants.file_limits import KNOWLEDGE_ALLOWED_EXTENSIONS
from app.services.document_intelligence.parsers import extract_document


def test_knowledge_allowlist_includes_pptx_and_md():
    assert KNOWLEDGE_ALLOWED_EXTENSIONS[".pptx"] == "pptx"
    assert KNOWLEDGE_ALLOWED_EXTENSIONS[".md"] == "md"
    assert KNOWLEDGE_ALLOWED_EXTENSIONS[".pdf"] == "pdf"


def test_extract_markdown(tmp_path: Path):
    path = tmp_path / "guide.md"
    path.write_text("# Wireless Design\n\nAP spacing guidance.\n", encoding="utf-8")
    result = extract_document(path, "md")
    assert "Wireless Design" in result.full_text
    assert result.metadata.get("parser") == "markdown"
    assert result.metadata.get("heading_1") == "Wireless Design"


def test_extract_pptx(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Cloud Reference Architecture"
    prs.save(str(path))

    result = extract_document(path, "pptx")
    assert "Cloud Reference Architecture" in result.full_text
    assert result.metadata.get("parser") == "python-pptx"
    assert len(result.pages) >= 1
